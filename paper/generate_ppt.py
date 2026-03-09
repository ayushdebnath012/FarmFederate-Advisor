"""
FarmFederate PPT Generator
Run: python generate_ppt.py
Requires: pip install python-pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Color palette ──────────────────────────────────────────────────────────────
GREEN_DARK  = RGBColor(0x1B, 0x5E, 0x20)   # deep green (title bar)
GREEN_MID   = RGBColor(0x2E, 0x7D, 0x32)   # medium green (accents)
GREEN_LIGHT = RGBColor(0xA5, 0xD6, 0xA7)   # light green (bg strip)
LLM_BLUE    = RGBColor(0x41, 0x69, 0xE1)
VIT_GREEN   = RGBColor(0x22, 0x8B, 0x22)
VLM_RED     = RGBColor(0xB2, 0x22, 0x22)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK        = RGBColor(0x21, 0x21, 0x21)
GRAY        = RGBColor(0x75, 0x75, 0x75)
YELLOW      = RGBColor(0xFF, 0xD6, 0x00)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank

# ── Helpers ────────────────────────────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.width = line_w
    if fill:
        shape.fill.solid(); shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT,
             wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb

def header_bar(slide, title, subtitle=""):
    add_rect(slide, 0, 0, 13.33, 1.2, fill=GREEN_DARK)
    add_text(slide, title, 0.3, 0.08, 12.5, 0.65,
             size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.3, 0.72, 12.5, 0.4,
                 size=14, color=GREEN_LIGHT, align=PP_ALIGN.LEFT)

def footer(slide, num):
    add_rect(slide, 0, 7.15, 13.33, 0.35, fill=GREEN_DARK)
    add_text(slide, "FarmFederate | Federated LLM · ViT · VLM for Crop Stress Detection",
             0.2, 7.15, 10, 0.32, size=9, color=GREEN_LIGHT)
    add_text(slide, str(num), 12.8, 7.15, 0.4, 0.32,
             size=9, color=WHITE, align=PP_ALIGN.RIGHT)

def bullet_box(slide, items, l, t, w, h, title="", color=DARK, size=15):
    if title:
        add_text(slide, title, l, t, w, 0.35, size=13, bold=True, color=GREEN_MID)
        t += 0.35
    for item in items:
        add_text(slide, f"▸  {item}", l, t, w, 0.4, size=size, color=color)
        t += 0.38

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xF1, 0xF8, 0xE9))

# big green band
add_rect(sl, 0, 1.8, 13.33, 3.5, fill=GREEN_DARK)

add_text(sl, "FarmFederate", 0.5, 2.0, 12.3, 1.1,
         size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl,
         "Federated LLM, ViT, and VLM for Privacy-Preserving Crop Stress Detection",
         0.5, 3.1, 12.3, 0.9, size=22, color=GREEN_LIGHT, align=PP_ALIGN.CENTER)

add_text(sl, "Ayush Debnath · Ruelia Saha · Sudip Misra",
         0.5, 4.2, 12.3, 0.5, size=16, color=WHITE, align=PP_ALIGN.CENTER, italic=True)

# three paradigm badges
for x, label, col in [(1.5, "Federated LLM", LLM_BLUE),
                       (5.2, "Federated ViT", VIT_GREEN),
                       (8.9, "Federated VLM", VLM_RED)]:
    add_rect(sl, x, 5.5, 3.0, 0.7, fill=col)
    add_text(sl, label, x, 5.52, 3.0, 0.65,
             size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

footer(sl, 1)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
header_bar(sl, "The Problem", "Why existing crop stress AI falls short")

problems = [
    "Crop stress (drought, disease, pests, heat) causes 20–40% yield loss globally",
    "Existing AI methods are CENTRALIZED — farms must share raw data",
    "Farms refuse to share data: privacy laws, competitive sensitivity, bandwidth",
    "Most methods are UNIMODAL — image-only or text-only, not both",
    "No prior work compares Federated LLM vs ViT vs VLM as independent paradigms",
]
bullet_box(sl, problems, 0.5, 1.4, 8.5, 5.0, size=16)

# example box on right
add_rect(sl, 9.3, 1.4, 3.7, 5.0, fill=RGBColor(0xFF, 0xF9, 0xC4), line=YELLOW, line_w=Pt(1.5))
add_text(sl, "Real-World Example", 9.4, 1.5, 3.5, 0.4, size=13, bold=True, color=DARK)
add_text(sl,
         "Farm A (Punjab) has drought data.\nFarm B (Kerala) has disease data.\nFarm C (UP) has pest data.\n\n"
         "Centralized: all must upload photos to one server.\n\n"
         "FarmFederate: only model weights travel — raw images stay on-farm.",
         9.4, 1.95, 3.5, 4.2, size=12, color=DARK)
footer(sl, 2)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — FarmFederate Overview
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
header_bar(sl, "FarmFederate Overview", "Three independent federated paradigms, one framework")

# Three column cards
for x, col, label, sub, f1, ratio in [
    (0.3,  LLM_BLUE, "Federated LLM",  "Text: symptom descriptions",    "0.548", "92.9%"),
    (4.6,  VIT_GREEN,"Federated ViT",  "Image: crop photographs",        "0.659", "86.1%"),
    (8.9,  VLM_RED,  "Federated VLM",  "Text + Image multimodal fusion", "0.785", "92.6%"),
]:
    add_rect(sl, x, 1.35, 3.9, 5.7, fill=RGBColor(0xFA,0xFA,0xFA), line=col, line_w=Pt(2))
    add_rect(sl, x, 1.35, 3.9, 0.55, fill=col)
    add_text(sl, label, x+0.1, 1.37, 3.7, 0.5, size=17, bold=True, color=WHITE)
    add_text(sl, sub,   x+0.1, 1.97, 3.7, 0.4, size=12, color=GRAY, italic=True)

    add_text(sl, "Federated F1",  x+0.2, 2.55, 1.8, 0.35, size=11, color=GRAY)
    add_text(sl, f1,              x+0.2, 2.85, 1.8, 0.6,  size=30, bold=True, color=col)
    add_text(sl, "Fed/Cent",      x+2.1, 2.55, 1.6, 0.35, size=11, color=GRAY)
    add_text(sl, ratio,           x+2.1, 2.85, 1.6, 0.6,  size=24, bold=True, color=DARK)

# Common arch info
add_text(sl, "Common: FedAvg · K=3 clients · T=8 rounds · E=3 local epochs · 5-class crop stress",
         0.3, 7.0, 12.5, 0.4, size=11, color=GRAY, italic=True)
footer(sl, 3)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Paradigm 1: Federated LLM
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
header_bar(sl, "Paradigm 1 — Federated LLM", "Distributed text classification of crop symptom descriptions")

# Left: model table
add_text(sl, "5 LLM Encoder Variants", 0.4, 1.35, 5.5, 0.4, size=14, bold=True, color=LLM_BLUE)
rows = [("Model","Cent. F1","Fed. F1"),
        ("DistilBERT","0.571","0.530"),
        ("BERT-tiny","0.585","0.543"),
        ("RoBERTa-tiny","0.585","0.541"),
        ("ALBERT-tiny ★","0.590","0.548"),
        ("MobileBERT","0.535","0.497")]
for i,(m,c,f) in enumerate(rows):
    bg = RGBColor(0xE8,0xF5,0xE9) if i==0 else (RGBColor(0xE3,0xF2,0xFD) if "★" in m else WHITE)
    add_rect(sl, 0.4, 1.8+i*0.52, 5.5, 0.50, fill=bg, line=RGBColor(0xDD,0xDD,0xDD), line_w=Pt(0.5))
    bld = (i==0)
    add_text(sl, m, 0.5,  1.83+i*0.52, 2.6, 0.44, size=12, bold=bld, color=DARK)
    add_text(sl, c, 3.1,  1.83+i*0.52, 1.3, 0.44, size=12, bold=bld, color=DARK, align=PP_ALIGN.CENTER)
    add_text(sl, f, 4.4,  1.83+i*0.52, 1.3, 0.44, size=12, bold=bld, color=LLM_BLUE, align=PP_ALIGN.CENTER)

# Right: key findings
add_text(sl, "Key Findings", 6.5, 1.35, 6.5, 0.4, size=14, bold=True, color=LLM_BLUE)
findings = [
    "Most federation-robust paradigm (7.1% accuracy drop)",
    "Text features degrade minimally under non-IID partitioning",
    "100% prediction diversity maintained across all rounds",
    "Ideal for bandwidth-limited or low-compute deployments",
]
bullet_box(sl, findings, 6.5, 1.8, 6.5, 3.0, size=14)

# Example
add_rect(sl, 6.5, 4.4, 6.5, 2.2, fill=RGBColor(0xE3,0xF2,0xFD), line=LLM_BLUE, line_w=Pt(1))
add_text(sl, "Example Prediction", 6.65, 4.48, 6.2, 0.38, size=12, bold=True, color=LLM_BLUE)
add_text(sl,
         'Input text:  "Brown spots on leaf margins, leaves curling upward, '
         'dry cracked soil near base of plant"\n\n'
         'Federated LLM prediction →  Disease Risk  (61% confidence)',
         6.65, 4.9, 6.2, 1.6, size=12, color=DARK)
footer(sl, 4)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Paradigm 2: Federated ViT
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
header_bar(sl, "Paradigm 2 — Federated ViT", "Distributed image classification of crop photographs")

rows = [("Model","Cent. F1","Fed. F1"),
        ("ViT-Base","0.696","0.599"),
        ("DeiT-tiny","0.747","0.643"),
        ("Swin-tiny","0.724","0.623"),
        ("ConvNeXT-tiny ★","0.765","0.659"),
        ("EfficientNet","0.724","0.623")]
for i,(m,c,f) in enumerate(rows):
    bg = RGBColor(0xE8,0xF5,0xE9) if i==0 else (RGBColor(0xF1,0xF8,0xE9) if "★" in m else WHITE)
    add_rect(sl, 0.4, 1.8+i*0.52, 5.5, 0.50, fill=bg, line=RGBColor(0xDD,0xDD,0xDD), line_w=Pt(0.5))
    bld = (i==0)
    add_text(sl, m, 0.5,  1.83+i*0.52, 2.7, 0.44, size=12, bold=bld, color=DARK)
    add_text(sl, c, 3.2,  1.83+i*0.52, 1.3, 0.44, size=12, bold=bld, color=DARK, align=PP_ALIGN.CENTER)
    add_text(sl, f, 4.5,  1.83+i*0.52, 1.3, 0.44, size=12, bold=bld, color=VIT_GREEN, align=PP_ALIGN.CENTER)

add_text(sl, "Key Findings", 6.5, 1.35, 6.5, 0.4, size=14, bold=True, color=VIT_GREEN)
findings = [
    "Highest unimodal accuracy (centralized F1 = 0.765)",
    "Largest federated gap (13.9%) — visual features are non-IID sensitive",
    "Residual CNN encoders with batch norm + spatial dropout",
    "Future fix: FedBN or pretrained ViT weights with LoRA fine-tuning",
]
bullet_box(sl, findings, 6.5, 1.8, 6.5, 3.0, size=14)

add_rect(sl, 6.5, 4.4, 6.5, 2.2, fill=RGBColor(0xF1,0xF8,0xE9), line=VIT_GREEN, line_w=Pt(1))
add_text(sl, "Example Prediction", 6.65, 4.48, 6.2, 0.38, size=12, bold=True, color=VIT_GREEN)
add_text(sl,
         'Input image:  224×224 crop photo — pale greenish-yellow leaves,\n'
         'drooping stems, vertical stripe pattern visible\n\n'
         'Federated ViT prediction →  Water Stress  (73% confidence)',
         6.65, 4.9, 6.2, 1.6, size=12, color=DARK)
footer(sl, 5)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Paradigm 3: Federated VLM
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
header_bar(sl, "Paradigm 3 — Federated VLM", "Distributed multimodal fusion of text + image")

fusions = [("Concatenation","0.797","0.738"),("Cross-Attention","0.807","0.747"),
           ("Gated","0.779","0.722"),("CLIP ★","0.848","0.785"),
           ("Flamingo","0.816","0.755"),("BLIP-2","0.834","0.771"),
           ("CoCa","0.783","0.725"),("Unified-IO","0.802","0.743")]
header_row = ("Fusion","Cent. F1","Fed. F1")
all_rows = [header_row] + fusions
for i,(m,c,f) in enumerate(all_rows):
    bg = RGBColor(0xFF,0xEB,0xEE) if i==0 else (RGBColor(0xFF,0xCD,0xD2) if "★" in m else WHITE)
    add_rect(sl, 0.3, 1.8+i*0.47, 5.8, 0.45, fill=bg, line=RGBColor(0xDD,0xDD,0xDD), line_w=Pt(0.5))
    bld = (i==0)
    add_text(sl, m, 0.4,  1.83+i*0.47, 2.8, 0.4, size=11, bold=bld, color=DARK)
    add_text(sl, c, 3.2,  1.83+i*0.47, 1.4, 0.4, size=11, bold=bld, color=DARK,    align=PP_ALIGN.CENTER)
    add_text(sl, f, 4.6,  1.83+i*0.47, 1.4, 0.4, size=11, bold=bld, color=VLM_RED, align=PP_ALIGN.CENTER)

add_text(sl, "Key Findings", 6.5, 1.35, 6.5, 0.4, size=14, bold=True, color=VLM_RED)
findings = [
    "Best absolute federated F1 = 0.785 (CLIP fusion)",
    "Fed/Cent ratio = 92.6% — comparable to Federated LLM",
    "CLIP contrastive alignment outperforms complex Q-Former (BLIP-2)",
    "Multimodal gain: +0.083 over best unimodal (ViT), +0.237 over LLM",
    "Recommended paradigm when compute is available",
]
bullet_box(sl, findings, 6.5, 1.8, 6.5, 3.0, size=13)

add_rect(sl, 6.5, 4.55, 6.5, 2.1, fill=RGBColor(0xFF,0xEB,0xEE), line=VLM_RED, line_w=Pt(1))
add_text(sl, "Example Prediction", 6.65, 4.63, 6.2, 0.38, size=12, bold=True, color=VLM_RED)
add_text(sl,
         'Image: small circular holes in leaves, sticky residue on stems\n'
         'Text:  "whitish deposits under leaves, plant stunted"\n\n'
         'Federated VLM (CLIP) →  Pest Risk  (84.8% confidence)',
         6.65, 5.07, 6.2, 1.45, size=12, color=DARK)
footer(sl, 6)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Cross-Paradigm Comparison
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
header_bar(sl, "Cross-Paradigm Comparison", "Which federated paradigm wins and when?")

# Big comparison table
cols = ["Paradigm","Cent. F1","Fed. F1","Fed/Cent","Gap","Verdict"]
data = [
    ("Federated LLM",  "0.590","0.548","92.9%","7.1%",  "Most federation-robust"),
    ("Federated ViT",  "0.765","0.659","86.1%","13.9%", "Best unimodal accuracy"),
    ("Federated VLM",  "0.848","0.785","92.6%","7.4%",  "Best overall ★"),
]
col_colors = [LLM_BLUE, VIT_GREEN, VLM_RED]

xs = [0.3, 2.3, 3.9, 5.5, 7.0, 8.3]
ws = [2.0, 1.6, 1.6, 1.5, 1.3, 4.7]
# header
for j,(col,x,w) in enumerate(zip(cols,xs,ws)):
    add_rect(sl, x, 1.35, w, 0.42, fill=GREEN_DARK)
    add_text(sl, col, x+0.05, 1.37, w-0.1, 0.38, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for i,(row, rc) in enumerate(zip(data, col_colors)):
    bg = RGBColor(0xF5,0xF5,0xF5) if i%2==0 else WHITE
    for j,(val,x,w) in enumerate(zip(row,xs,ws)):
        add_rect(sl, x, 1.78+i*0.58, w, 0.55, fill=bg, line=RGBColor(0xCC,0xCC,0xCC), line_w=Pt(0.5))
        col_c = rc if j in (2,3) else DARK
        bld   = (j==5)
        add_text(sl, val, x+0.05, 1.82+i*0.58, w-0.1, 0.48,
                 size=13, bold=bld, color=col_c, align=PP_ALIGN.CENTER)

# Decision guide
add_rect(sl, 0.3, 3.75, 12.7, 2.9, fill=RGBColor(0xF9,0xFB,0xE7), line=GREEN_MID, line_w=Pt(1.5))
add_text(sl, "Decision Guide", 0.5, 3.83, 5.0, 0.4, size=14, bold=True, color=GREEN_MID)
guide = [
    "Use Federated LLM  →  bandwidth-constrained farms, text-rich sensor logs",
    "Use Federated ViT  →  image-heavy setups, future FedBN aggregation",
    "Use Federated VLM  →  when compute & dual-modality data are available (recommended)",
]
bullet_box(sl, guide, 0.5, 4.28, 12.2, 2.2, size=13)
footer(sl, 7)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Dataset & Anti-Collapse
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
header_bar(sl, "Dataset & Anti-Collapse Mechanisms", "Real data + robust training for all three paradigms")

# Left — dataset
add_text(sl, "5-Class Crop Stress Dataset", 0.4, 1.35, 6.2, 0.4, size=14, bold=True, color=GREEN_MID)
classes = [("Water Stress","Drought, wilting, stomatal closure"),
           ("Nutrient Def.","Chlorosis, stunted growth"),
           ("Pest Risk","Holes, sticky residue, webbing"),
           ("Disease Risk","Lesions, spots, fungal/bacterial"),
           ("Heat Stress","Leaf scorch, tip burn")]
for i,(c,d) in enumerate(classes):
    add_rect(sl, 0.4, 1.82+i*0.62, 6.2, 0.58, fill=RGBColor(0xE8,0xF5,0xE9),
             line=GREEN_MID, line_w=Pt(0.5))
    add_text(sl, c, 0.55, 1.85+i*0.62, 2.0, 0.52, size=12, bold=True, color=GREEN_DARK)
    add_text(sl, d, 2.55, 1.85+i*0.62, 3.9, 0.52, size=11, color=GRAY)

add_text(sl, "Text: AG News + PubMed + SQuAD → ~1,089 balanced samples\n"
             "Images: PlantVillage + Beans (real) + Synthetic (tuned difficulty)\n"
             "Original imbalance: 25:1 (disease:heat) → rebalanced via capping",
         0.4, 5.15, 6.2, 1.0, size=11, color=GRAY, italic=True)

# Right — anti-collapse
add_text(sl, "Anti-Collapse Mechanisms", 7.0, 1.35, 5.9, 0.4, size=14, bold=True, color=VLM_RED)
mechanisms = [
    ("Balanced Batch Sampling", "Equal class samples per batch via oversampling minorities"),
    ("Diversity Loss",          "Penalizes low entropy: λ=1.0, 70% threshold (Eq. 3)"),
    ("Early Collapse Abort",    "Stops after 3 consecutive collapsed epochs, restores best ckpt"),
    ("Focal Loss",              "Down-weights easy majority class (γ=2.0, sqrt class weights)"),
]
for i,(title,desc) in enumerate(mechanisms):
    add_rect(sl, 7.0, 1.82+i*1.1, 5.9, 1.05, fill=RGBColor(0xFF,0xEB,0xEE),
             line=VLM_RED, line_w=Pt(0.5))
    add_text(sl, title, 7.15, 1.85+i*1.1, 5.6, 0.42, size=12, bold=True, color=VLM_RED)
    add_text(sl, desc,  7.15, 2.27+i*1.1, 5.6, 0.55, size=11, color=DARK)

add_text(sl, "Result: 100% prediction diversity maintained across all 24 trained models",
         7.0, 6.25, 5.9, 0.45, size=12, bold=True, color=GREEN_DARK)
footer(sl, 8)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Results Summary
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
header_bar(sl, "Results Summary", "Key numbers across all three federated paradigms")

metrics = [
    ("Best Federated VLM F1",  "0.785", VLM_RED),
    ("Best Centralized VLM F1","0.848", VLM_RED),
    ("Avg Fed/Cent Retention",  "90.5%", GREEN_MID),
    ("Stress-biased Eval F1",   "≥0.989",GREEN_DARK),
    ("Prediction Diversity",    "100%",  GREEN_MID),
    ("Total Models Trained",    "24",    GRAY),
]
for i,(label,val,col) in enumerate(metrics):
    x = 0.4 + (i%3)*4.3
    y = 1.35 + (i//3)*2.6
    add_rect(sl, x, y, 4.0, 2.3, fill=RGBColor(0xFA,0xFA,0xFA),
             line=col, line_w=Pt(2))
    add_text(sl, val,   x+0.15, y+0.3,  3.7, 1.1, size=36, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(sl, label, x+0.15, y+1.45, 3.7, 0.7, size=12, color=DARK, align=PP_ALIGN.CENTER)

footer(sl, 9)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Practical Example (End-to-End)
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
header_bar(sl, "End-to-End Example", "A farmer uses FarmFederate on Android")

steps = [
    ("1. Capture", "Farmer photographs yellowing maize leaves using the FarmFederate app"),
    ("2. Describe", 'Types: "leaves pale yellow, lower leaves affected first, slow growth"'),
    ("3. Federated VLM", "App sends encoded features (NOT raw photo) to federated server via CLIP fusion"),
    ("4. FedAvg Round", "Server aggregates updates from Farm A (Punjab), B (Kerala), C (UP)"),
    ("5. Prediction", "Nutrient Deficiency — 78% confidence  (Nitrogen deficiency pattern)"),
    ("6. Advice", "Apply 30 kg urea/acre. Re-evaluate in 10 days. Privacy: photo never left device."),
]
for i,(title,desc) in enumerate(steps):
    row, col_idx = divmod(i, 3)
    x = 0.3 + col_idx * 4.35
    y = 1.35 + row * 2.7
    col = [LLM_BLUE, GREEN_MID, VLM_RED][col_idx]
    add_rect(sl, x, y, 4.1, 2.55, fill=RGBColor(0xFA,0xFA,0xFA), line=col, line_w=Pt(1.5))
    add_rect(sl, x, y, 4.1, 0.5, fill=col)
    add_text(sl, title, x+0.1, y+0.05, 3.9, 0.42, size=14, bold=True, color=WHITE)
    add_text(sl, desc,  x+0.1, y+0.6,  3.9, 1.85, size=12, color=DARK)

footer(sl, 10)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Conclusion
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xF1,0xF8,0xE9))
header_bar(sl, "Conclusion", "Three paradigms, one answer")

conclusions = [
    ("Federated LLM (ALBERT-tiny)",
     "F1=0.548, 92.9% retention — most federation-robust. Use when bandwidth is scarce.",
     LLM_BLUE),
    ("Federated ViT (ConvNeXT-tiny)",
     "F1=0.659, 86.1% retention — best unimodal F1 but most sensitive to non-IID heterogeneity.",
     VIT_GREEN),
    ("Federated VLM (CLIP fusion)",
     "F1=0.785, 92.6% retention — best overall. Recommended when dual-modality data is available.",
     VLM_RED),
    ("Anti-Collapse Works",
     "100% prediction diversity across all 24 models via diversity loss + balanced sampling.",
     GREEN_DARK),
]
for i,(title,desc,col) in enumerate(conclusions):
    add_rect(sl, 0.4, 1.4+i*1.3, 12.5, 1.2, fill=WHITE, line=col, line_w=Pt(2))
    add_rect(sl, 0.4, 1.4+i*1.3, 0.25, 1.2, fill=col)
    add_text(sl, title, 0.8, 1.45+i*1.3, 5.5, 0.42, size=13, bold=True, color=col)
    add_text(sl, desc,  0.8, 1.87+i*1.3, 12.0, 0.55, size=12, color=DARK)

add_text(sl, "FarmFederate is the first work to define, independently train, and directly compare "
             "Federated LLM, ViT, and VLM as distinct paradigms for agricultural federated AI.",
         0.4, 6.65, 12.5, 0.55, size=12, bold=True, color=GREEN_DARK, italic=True)
footer(sl, 11)

# ── Save ───────────────────────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(__file__), "FarmFederate_Presentation.pptx")
prs.save(out)
print(f"Saved: {out}")
print("  11 slides generated.")
