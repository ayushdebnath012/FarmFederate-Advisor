from __future__ import annotations

import datetime as dt
import shutil
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "FarmFederate_ppt_v2.pptx"
BACKUP = ROOT / f"FarmFederate_ppt_v2_before_paper_update_{dt.datetime.now():%Y%m%d_%H%M%S}.pptx"
PREVIEW_DIR = ROOT / "ppt_previews_v2_updated"


W, H = Inches(13.333), Inches(7.5)

INK = RGBColor(11, 16, 32)
INK2 = RGBColor(24, 34, 56)
PAPER = RGBColor(251, 250, 247)
MIST = RGBColor(236, 246, 240)
GREEN = RGBColor(27, 107, 67)
GREEN2 = RGBColor(48, 148, 96)
TEAL = RGBColor(14, 116, 144)
BLUE = RGBColor(37, 99, 235)
RED = RGBColor(198, 40, 40)
AMBER = RGBColor(224, 138, 0)
GRAY = RGBColor(102, 112, 133)
LIGHT = RGBColor(232, 238, 234)
WHITE = RGBColor(255, 255, 255)

FONT_HEAD = "Aptos Display"
FONT_BODY = "Aptos"


def rgb_hex(rgb: RGBColor) -> str:
    return f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"


def set_fill(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def no_line(shape):
    shape.line.fill.background()


def add_bg(slide, color=PAPER):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    set_fill(bg, color)
    no_line(bg)
    return bg


def add_band(slide, x, y, w, h, color, alpha=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    set_fill(shp, color)
    if alpha is not None:
        shp.fill.transparency = alpha
    no_line(shp)
    return shp


def add_text(slide, text, x, y, w, h, size=24, color=INK, bold=False,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, font=FONT_BODY,
             italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return tb


def add_title(slide, title, subtitle=None, kicker=None):
    if kicker:
        add_text(slide, kicker.upper(), Inches(0.55), Inches(0.33), Inches(4.3), Inches(0.25),
                 9, GREEN, True)
    add_text(slide, title, Inches(0.55), Inches(0.55), Inches(8.6), Inches(0.58),
             28, INK, True, font=FONT_HEAD)
    if subtitle:
        add_text(slide, subtitle, Inches(0.57), Inches(1.13), Inches(9.4), Inches(0.36),
                 12.5, GRAY)


def add_footer(slide, num, total=16):
    add_text(slide, "FarmFederate | paper-aligned deck", Inches(0.55), Inches(7.13),
             Inches(3.6), Inches(0.18), 7.5, GRAY)
    add_text(slide, f"{num} / {total}", Inches(12.15), Inches(7.12),
             Inches(0.6), Inches(0.18), 8, GRAY, align=PP_ALIGN.RIGHT)


def add_chip(slide, text, x, y, w, color=GREEN, text_color=WHITE):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(0.33))
    set_fill(shp, color)
    no_line(shp)
    add_text(slide, text, x + Inches(0.09), y + Inches(0.065), w - Inches(0.18),
             Inches(0.15), 8.8, text_color, True, align=PP_ALIGN.CENTER)
    return shp


def add_metric(slide, value, label, x, y, w, h, color=GREEN, note=None):
    add_text(slide, value, x, y, w, Inches(0.48), 25, color, True, align=PP_ALIGN.CENTER,
             font=FONT_HEAD)
    add_text(slide, label, x + Inches(0.02), y + Inches(0.52), w - Inches(0.04),
             Inches(0.36), 10.2, INK, True, align=PP_ALIGN.CENTER)
    if note:
        add_text(slide, note, x + Inches(0.05), y + Inches(0.91), w - Inches(0.1),
                 Inches(0.3), 8.6, GRAY, align=PP_ALIGN.CENTER)


def add_callout(slide, title, body, x, y, w, h, color=GREEN):
    add_band(slide, x, y, Inches(0.06), h, color)
    add_text(slide, title, x + Inches(0.18), y, w - Inches(0.18), Inches(0.25),
             13, color, True)
    add_text(slide, body, x + Inches(0.18), y + Inches(0.32), w - Inches(0.18),
             h - Inches(0.32), 11, INK)


def add_image_contain(slide, path, x, y, w, h, border=True, bg=WHITE):
    path = Path(path)
    if border:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        set_fill(box, bg)
        box.line.color.rgb = LIGHT
        box.line.width = Pt(1)
    if not path.exists():
        add_text(slide, f"Missing image: {path.name}", x + Inches(0.1), y + Inches(0.1),
                 w - Inches(0.2), h - Inches(0.2), 10, RED)
        return None
    with Image.open(path) as im:
        iw, ih = im.size
    ratio = min(w / iw, h / ih)
    pw, ph = int(iw * ratio), int(ih * ratio)
    px = x + (w - pw) / 2
    py = y + (h - ph) / 2
    return slide.shapes.add_picture(str(path), px, py, width=pw, height=ph)


def add_small_table(slide, rows, x, y, w, row_h=0.34, col_fr=(0.42, 0.27, 0.31),
                    header=True):
    cols = [w * f for f in col_fr]
    cy = y
    for r, row in enumerate(rows):
        if header and r == 0:
            add_band(slide, x, cy, w, Inches(row_h), INK2)
            ccol = WHITE
            bold = True
        else:
            if r % 2 == 0:
                add_band(slide, x, cy, w, Inches(row_h), MIST)
            ccol = INK
            bold = r == 0
        cx = x
        for c, cell in enumerate(row):
            add_text(slide, str(cell), cx + Inches(0.06), cy + Inches(0.075),
                     cols[c] - Inches(0.08), Inches(row_h - 0.05), 8.6, ccol, bold)
            cx += cols[c]
        cy += Inches(row_h)


def add_flow_node(slide, text, x, y, w, h, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    set_fill(shp, color)
    no_line(shp)
    add_text(slide, text, x + Inches(0.12), y + Inches(0.12), w - Inches(0.24),
             h - Inches(0.18), 11, WHITE, True, align=PP_ALIGN.CENTER,
             valign=MSO_ANCHOR.MIDDLE)
    return shp


def add_arrow(slide, x1, y1, x2, y2, color=GRAY):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(2)
    line.line.end_arrowhead = True
    return line


def new_slide(prs, idx, title, subtitle=None, kicker=None, bg=PAPER):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, bg)
    add_title(slide, title, subtitle, kicker)
    add_footer(slide, idx)
    return slide


def build_deck():
    if OUT.exists():
        shutil.copy2(OUT, BACKUP)

    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    slides_for_preview = []

    # 1. Cover
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, INK)
    add_band(s, 0, 0, Inches(13.333), Inches(0.18), GREEN2)
    add_image_contain(s, ROOT / "disease_annotated" / "_grid_overview.jpg",
                      Inches(8.3), Inches(0.45), Inches(4.45), Inches(6.55),
                      border=False)
    add_band(s, Inches(7.85), 0, Inches(0.18), H, GREEN2)
    add_text(s, "FarmFederate", Inches(0.72), Inches(1.55), Inches(6.4), Inches(0.72),
             44, WHITE, True, font=FONT_HEAD)
    add_text(s, "Multimodal federated learning for tea leaf disease detection",
             Inches(0.78), Inches(2.35), Inches(6.25), Inches(0.7), 18, RGBColor(217, 232, 224))
    add_text(s, "Final paper-aligned presentation", Inches(0.8), Inches(0.72),
             Inches(4.2), Inches(0.25), 11, GREEN2, True)
    for i, (v, lab) in enumerate([("5", "tea diseases"), ("18", "model variants"),
                                  ("0.949", "best macro F1"), ("0", "Kaggle images used")]):
        add_metric(s, v, lab, Inches(0.75 + i * 1.78), Inches(4.9), Inches(1.55),
                   Inches(1.15), [GREEN2, TEAL, AMBER, RED][i])
    add_text(s, "Field images + class-balanced augmentation + generated symptom text",
             Inches(0.8), Inches(6.55), Inches(6.8), Inches(0.28), 11, RGBColor(199, 213, 206))
    slides_for_preview.append(("FarmFederate", ["5 diseases", "18 models", "F1 0.949", "No Kaggle image corpus"], [ROOT / "disease_annotated" / "_grid_overview.jpg"]))

    # 2. Agenda
    s = new_slide(prs, 2, "Agenda", "What changed in the paper is now reflected in the deck.", "Roadmap")
    agenda = [
        ("01", "Problem and disease scope"),
        ("02", "Dataset, labels, and no-Kaggle training claim"),
        ("03", "Architecture: LLM + ViT + VLM + FedAvg + RAG"),
        ("04", "Training safeguards and non-IID simulation"),
        ("05", "Results: baselines, VLM fusion, federated retention"),
        ("06", "Visual fallback, SOTA benchmark, limitations, next steps"),
    ]
    for i, (num, txt) in enumerate(agenda):
        y = Inches(1.85 + i * 0.66)
        add_text(s, num, Inches(0.9), y, Inches(0.48), Inches(0.25), 13, GREEN, True)
        add_text(s, txt, Inches(1.45), y - Inches(0.03), Inches(8.8), Inches(0.3), 17, INK, True)
    add_chip(s, "Updated from FarmFederate_Paper_v3", Inches(9.55), Inches(6.62), Inches(2.75), GREEN)
    slides_for_preview.append(("Agenda", [t for _, t in agenda], []))

    # 3. Problem and disease scope
    s = new_slide(prs, 3, "Five tea diseases, one field diagnosis problem",
                  "The deck now names the disease scope and the visual symptom families.", "Problem")
    disease_rows = [
        ("Disease", "Real imgs", "Short field cue"),
        ("Leaf blight", "35", "Brown water-soaked necrosis"),
        ("Leaf hoppers", "8", "Puncture lesions and hopper damage"),
        ("Leaf rust", "62", "Orange-yellow rust pustules"),
        ("Looper caterpillars", "48", "Ragged holes, skeletonised patches"),
        ("Mosquito bug", "47", "Corky wounds and shoot-tip dieback"),
    ]
    add_small_table(s, disease_rows, Inches(0.75), Inches(1.78), Inches(6.25), row_h=0.47,
                    col_fr=(0.40, 0.18, 0.42))
    add_callout(s, "Why it is hard",
                "The classes mix fungal spots, rust pustules, and insect feeding damage. The model must separate texture cues from pest-injury patterns, not just green-leaf backgrounds.",
                Inches(7.35), Inches(1.75), Inches(4.85), Inches(1.55), GREEN)
    thumb_paths = [
        ROOT / "plots" / "real_leaf_blight.jpg",
        ROOT / "plots" / "real_leaf_hoppers.jpg",
        ROOT / "plots" / "real_leaf_rust.jpg",
        ROOT / "plots" / "real_looper_caterpillars.jpg",
        ROOT / "plots" / "real_mosquito_bug.jpg",
    ]
    labels = ["Blight", "Hoppers", "Rust", "Looper", "Mosquito"]
    for i, p in enumerate(thumb_paths):
        x = Inches(0.85 + i * 2.38)
        add_image_contain(s, p, x, Inches(5.0), Inches(1.75), Inches(1.1), border=True)
        add_text(s, labels[i], x, Inches(6.18), Inches(1.75), Inches(0.22), 9.5, INK,
                 True, align=PP_ALIGN.CENTER)
    slides_for_preview.append(("Five disease scope", [r[0] for r in disease_rows[1:]], thumb_paths[:2]))

    # 4. Dataset claim
    s = new_slide(prs, 4, "Dataset claim is now explicit and paper-aligned",
                  "Reported training uses field images, minority synthetic fill, and generated symptom text.", "Data")
    metrics = [("200", "field photos"), ("371", "YOLO-OBB boxes"), ("792", "training-run image tensors"), ("3,000", "generated text rows")]
    for i, (v, lab) in enumerate(metrics):
        add_metric(s, v, lab, Inches(0.78 + i * 3.05), Inches(1.85), Inches(2.35),
                   Inches(1.25), [GREEN, TEAL, AMBER, BLUE][i])
    add_callout(s, "No external Kaggle image corpus",
                "Kaggle credentials may be written by the Colab setup, but the reported training pipeline does not include a Kaggle image dataset.",
                Inches(0.9), Inches(3.7), Inches(5.65), Inches(1.25), RED)
    add_callout(s, "What was balanced",
                "The export prepares 800 images per class for inspection/download. The reported run loads the 200 sorted field images and adds minority fill only for leaf hoppers.",
                Inches(6.95), Inches(3.7), Inches(5.45), Inches(1.25), GREEN)
    add_image_contain(s, ROOT / "plots" / "plot21_dataset_comparison.png",
                      Inches(7.0), Inches(5.22), Inches(5.2), Inches(1.45), border=True)
    slides_for_preview.append(("Dataset claim", [m[0] + " " + m[1] for m in metrics], [ROOT / "plots" / "plot21_dataset_comparison.png"]))

    # 5. Data pipeline
    s = new_slide(prs, 5, "Colab data pipeline: field-first, then controlled augmentation",
                  "The paper now distinguishes exported data preparation from the reported training run.", "Pipeline")
    nodes = [
        ("Raw field collection\nIIT Kharagpur tea garden", GREEN),
        ("Sorted 5-class folders\n200 real images", TEAL),
        ("YOLO-OBB labels\n371 disease boxes", BLUE),
        ("Minority fill\n+592 leaf-hoppers", AMBER),
        ("Generated symptom text\n600 rows/class", RED),
    ]
    xs = [0.65, 3.05, 5.45, 7.85, 10.25]
    for i, (txt, col) in enumerate(nodes):
        add_flow_node(s, txt, Inches(xs[i]), Inches(2.05), Inches(1.8), Inches(0.86), col)
        if i < len(nodes) - 1:
            add_arrow(s, Inches(xs[i] + 1.82), Inches(2.48), Inches(xs[i + 1] - 0.04), Inches(2.48))
    add_text(s, "Training loaders", Inches(0.85), Inches(4.0), Inches(2.4), Inches(0.3),
             16, GREEN, True)
    loader_rows = [
        ("Text", "3,000 rows", "2400 train / 300 val"),
        ("Image", "792 tensors", "633 train / 79 val"),
        ("VLM pairs", "633 train", "79 validation pairs"),
    ]
    for i, row in enumerate(loader_rows):
        y = Inches(4.48 + i * 0.55)
        add_text(s, row[0], Inches(1.0), y, Inches(1.0), Inches(0.22), 11, INK, True)
        add_text(s, row[1], Inches(2.25), y, Inches(1.35), Inches(0.22), 11, BLUE, True)
        add_text(s, row[2], Inches(3.85), y, Inches(2.3), Inches(0.22), 11, GRAY)
    add_image_contain(s, ROOT / "plots" / "plot26_stress_distribution.png",
                      Inches(7.0), Inches(4.0), Inches(5.0), Inches(2.35), border=True)
    slides_for_preview.append(("Data pipeline", [n[0].replace("\n", " ") for n, _ in nodes], [ROOT / "plots" / "plot26_stress_distribution.png"]))

    # 6. Architecture
    s = new_slide(prs, 6, "Architecture: local multimodal learning plus federated aggregation",
                  "Text, image, fusion, advice retrieval, and FedAvg are separated cleanly.", "System")
    add_flow_node(s, "Symptom text\nLLM encoder", Inches(0.85), Inches(2.1), Inches(2.0), Inches(0.92), TEAL)
    add_flow_node(s, "Leaf image\nViT encoder", Inches(0.85), Inches(4.0), Inches(2.0), Inches(0.92), GREEN)
    add_flow_node(s, "8 VLM fusion\nstrategies", Inches(4.0), Inches(3.05), Inches(2.15), Inches(1.0), RED)
    add_flow_node(s, "Disease\nclassifier", Inches(7.15), Inches(3.05), Inches(1.85), Inches(1.0), BLUE)
    add_flow_node(s, "Local FAISS\nRAG advice", Inches(10.0), Inches(2.1), Inches(2.1), Inches(0.92), AMBER)
    add_flow_node(s, "FedAvg server\nweights only", Inches(10.0), Inches(4.0), Inches(2.1), Inches(0.92), INK2)
    for (x1, y1, x2, y2) in [(2.85, 2.56, 4.0, 3.3), (2.85, 4.46, 4.0, 3.8),
                             (6.15, 3.55, 7.15, 3.55), (9.0, 3.35, 10.0, 2.56),
                             (9.0, 3.75, 10.0, 4.46)]:
        add_arrow(s, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    add_text(s, "Privacy boundary: raw leaf images, field logs, and knowledge-base documents stay on the garden device.",
             Inches(1.0), Inches(5.7), Inches(10.6), Inches(0.35), 15, RED, True, align=PP_ALIGN.CENTER)
    slides_for_preview.append(("Architecture", ["LLM", "ViT", "8 VLM fusions", "RAG", "FedAvg"], []))

    # 7. Training safeguards
    s = new_slide(prs, 7, "Training safeguards prevent class collapse",
                  "The paper now explains why all five classes remain represented despite imbalance.", "Training")
    safeguards = [
        ("Focal loss", "Pushes learning toward harder and rarer examples."),
        ("Diversity regularization", "Entropy-based term discourages one-class prediction collapse."),
        ("Dampened class weights", "Square-root class weights capped at 10x."),
        ("Mixed precision", "T4-friendly training with AMP."),
        ("Non-IID clients", "Three simulated estate clients test privacy-utility retention."),
    ]
    for i, (head, body) in enumerate(safeguards):
        y = Inches(1.75 + i * 0.78)
        add_text(s, head, Inches(0.9), y, Inches(2.5), Inches(0.24), 14, GREEN, True)
        add_text(s, body, Inches(3.25), y, Inches(5.8), Inches(0.28), 12.5, INK)
    add_image_contain(s, ROOT / "plots" / "plot08_params.png",
                      Inches(9.25), Inches(1.8), Inches(3.0), Inches(3.25), border=True)
    add_text(s, "Lightweight models keep the Colab and edge-device story credible.",
             Inches(9.35), Inches(5.35), Inches(2.9), Inches(0.62), 11, GRAY, align=PP_ALIGN.CENTER)
    slides_for_preview.append(("Training safeguards", [h for h, _ in safeguards], [ROOT / "plots" / "plot08_params.png"]))

    # 8. Results overview
    s = new_slide(prs, 8, "Main result: multimodal fusion is the top-line winner",
                  "CLIP-style VLM fusion reaches F1=0.949 across 18 tested architectures.", "Results")
    add_image_contain(s, ROOT / "plots" / "plot04_model_type_overview.png",
                      Inches(0.8), Inches(1.55), Inches(7.3), Inches(4.95), border=True)
    add_metric(s, "0.949", "best VLM: CLIP", Inches(8.55), Inches(1.7), Inches(3.2), Inches(1.2), RED)
    add_metric(s, "0.911", "best image-only: EfficientNet", Inches(8.55), Inches(3.2), Inches(3.2), Inches(1.2), GREEN)
    add_metric(s, "0.487", "best text-only: BERT-tiny", Inches(8.55), Inches(4.7), Inches(3.2), Inches(1.2), BLUE)
    add_text(s, "VLM improves over best ViT by 0.038 F1 and over best LLM by 0.462 F1.",
             Inches(0.9), Inches(6.65), Inches(10.9), Inches(0.28), 13, INK, True, align=PP_ALIGN.CENTER)
    slides_for_preview.append(("Main result", ["VLM 0.949", "ViT 0.911", "LLM 0.487"], [ROOT / "plots" / "plot04_model_type_overview.png"]))

    # 9. Encoder baselines
    s = new_slide(prs, 9, "Encoder baselines: text helps, images carry visual discrimination",
                  "The new paper figures use larger labels and publication-scaled plots.", "Baselines")
    add_image_contain(s, ROOT / "plots" / "plot01_llm_comparison.png",
                      Inches(0.65), Inches(1.65), Inches(3.8), Inches(2.2), border=True)
    add_image_contain(s, ROOT / "plots" / "plot02_vit_comparison.png",
                      Inches(4.75), Inches(1.65), Inches(3.8), Inches(2.2), border=True)
    add_image_contain(s, ROOT / "plots" / "plot03_vlm_fusion_comparison.png",
                      Inches(8.85), Inches(1.65), Inches(3.8), Inches(2.2), border=True)
    add_callout(s, "Readout",
                "LLM scores remain limited because symptom templates are deliberately mixed. Vision models are much stronger, and VLM fusion adds the final gain when text and images are both present.",
                Inches(1.05), Inches(4.55), Inches(11.0), Inches(1.3), GREEN)
    slides_for_preview.append(("Encoder baselines", ["LLM", "ViT", "VLM"], [ROOT / "plots" / "plot03_vlm_fusion_comparison.png"]))

    # 10. VLM fusion ranking
    s = new_slide(prs, 10, "Eight VLM fusions were benchmarked under the same protocol",
                  "The deck now carries the full fusion comparison instead of only the winner.", "VLM")
    add_image_contain(s, ROOT / "plots" / "plot03_vlm_fusion_comparison.png",
                      Inches(0.8), Inches(1.55), Inches(5.35), Inches(4.95), border=True)
    fusion_rows = [
        ("Rank", "Fusion", "Macro F1"),
        ("1", "CLIP", "0.949"),
        ("2", "BLIP-2", "0.937"),
        ("3", "Cross-attn", "0.924"),
        ("4", "Flamingo / Unified-IO", "0.911"),
        ("8", "Concat", "0.886"),
    ]
    add_small_table(s, fusion_rows, Inches(6.75), Inches(1.85), Inches(5.3), row_h=0.46,
                    col_fr=(0.18, 0.52, 0.30))
    add_text(s, "The VLM spread is narrow relative to LLM-VLM, so the result is a modality gain, not one unstable architecture.",
             Inches(6.9), Inches(5.25), Inches(4.9), Inches(0.78), 13, INK, True)
    slides_for_preview.append(("VLM fusion", [r[1] for r in fusion_rows[1:]], [ROOT / "plots" / "plot03_vlm_fusion_comparison.png"]))

    # 11. Federated retention
    s = new_slide(prs, 11, "Federated learning preserves most of the multimodal gain",
                  "Raw data stays local; only model updates are aggregated.", "Federated")
    add_image_contain(s, ROOT / "plots" / "plot05_centralized_vs_federated.png",
                      Inches(0.75), Inches(1.55), Inches(5.65), Inches(3.1), border=True)
    add_image_contain(s, ROOT / "plots" / "plot27_federated_convergence.png",
                      Inches(6.85), Inches(1.55), Inches(5.65), Inches(3.1), border=True)
    add_metric(s, "107.7%", "LLM retention", Inches(1.2), Inches(5.2), Inches(2.3), Inches(1.0), BLUE)
    add_metric(s, "100.0%", "ViT retention", Inches(5.15), Inches(5.2), Inches(2.3), Inches(1.0), GREEN)
    add_metric(s, "98.6%", "VLM retention", Inches(9.1), Inches(5.2), Inches(2.3), Inches(1.0), RED)
    slides_for_preview.append(("Federated retention", ["LLM 107.7%", "ViT 100.0%", "VLM 98.6%"], [ROOT / "plots" / "plot05_centralized_vs_federated.png"]))

    # 12. RAG advisory
    s = new_slide(prs, 12, "Classify-Retrieve-Advise keeps the advisory layer local",
                  "The RAG module retrieves disease-specific treatment guidance without uploading farm records.", "RAG")
    add_image_contain(s, ROOT / "plots" / "rag_01_retrieval_scores.png",
                      Inches(0.7), Inches(1.55), Inches(5.6), Inches(3.1), border=True)
    add_image_contain(s, ROOT / "plots" / "rag_02_score_heatmap.png",
                      Inches(6.75), Inches(1.55), Inches(5.6), Inches(3.1), border=True)
    add_callout(s, "Paper update",
                "The retrieval heatmap confirms correct top-1 disease retrieval for all five classes; local FAISS stores one treatment document per disease.",
                Inches(1.0), Inches(5.25), Inches(10.9), Inches(0.95), TEAL)
    slides_for_preview.append(("RAG advisory", ["Top-1 retrieval all five classes", "Local FAISS"], [ROOT / "plots" / "rag_02_score_heatmap.png"]))

    # 13. Visual fallback
    s = new_slide(prs, 13, "Visual diagnostic fallback gives field users a concrete reference",
                  "Text-only uncertainty maps back to YOLO-OBB annotated leaf images.", "Fallback")
    add_image_contain(s, ROOT / "disease_annotated" / "_grid_overview.jpg",
                      Inches(0.9), Inches(1.35), Inches(6.0), Inches(5.35), border=True)
    add_callout(s, "What the user sees",
                "Representative leaf photos are overlaid with class-specific OBB polygons and disease tags. If annotated field examples are scarce, the visualizer can fall back to the synthetic tea_data pool.",
                Inches(7.35), Inches(1.65), Inches(4.7), Inches(1.65), GREEN)
    add_callout(s, "Why it belongs in the deck",
                "It turns the classifier into an operational diagnostic tool: farmers get both a label and a visual explanation for the likely disease region.",
                Inches(7.35), Inches(4.0), Inches(4.7), Inches(1.35), TEAL)
    slides_for_preview.append(("Visual fallback", ["OBB annotated grid", "Visual explanation"], [ROOT / "disease_annotated" / "_grid_overview.jpg"]))

    # 14. Heatmap and ablation
    s = new_slide(prs, 14, "Heatmaps and ablation show where the gain comes from",
                  "The paper now places result evidence before conclusion, not inside it.", "Evidence")
    add_image_contain(s, ROOT / "plots" / "plot13_heatmap.png",
                      Inches(0.65), Inches(1.45), Inches(6.2), Inches(4.55), border=True)
    add_image_contain(s, ROOT / "plots" / "plot10d_modality_contribution.png",
                      Inches(7.15), Inches(1.55), Inches(5.0), Inches(2.75), border=True)
    ablation = [
        ("Mode", "Best model", "F1"),
        ("Text only", "BERT-tiny", "0.487"),
        ("Image only", "EfficientNet", "0.911"),
        ("Multimodal", "VLM-CLIP", "0.949"),
    ]
    add_small_table(s, ablation, Inches(7.25), Inches(4.75), Inches(4.8), row_h=0.42,
                    col_fr=(0.35, 0.40, 0.25))
    slides_for_preview.append(("Heatmap + ablation", ["Text 0.487", "Image 0.911", "VLM 0.949"], [ROOT / "plots" / "plot13_heatmap.png"]))

    # 15. SOTA benchmark
    s = new_slide(prs, 15, "Benchmark placement: competitive with current tea disease AI",
                  "FarmFederate sits in the leading group while adding privacy-preserving multimodal training.", "Benchmark")
    add_image_contain(s, ROOT / "plots" / "plot11_paper_comparison.png",
                      Inches(0.75), Inches(1.35), Inches(7.0), Inches(5.2), border=True)
    add_metric(s, "#12", "rank out of 35 models", Inches(8.25), Inches(1.55), Inches(3.4), Inches(1.2), GREEN)
    add_metric(s, "0.949", "FarmFederate best F1", Inches(8.25), Inches(3.05), Inches(3.4), Inches(1.2), RED)
    add_metric(s, "0.892", "mean SOTA F1", Inches(8.25), Inches(4.55), Inches(3.4), Inches(1.2), BLUE)
    add_text(s, "Competitive top-line performance plus local-data privacy is the main contribution.",
             Inches(8.1), Inches(6.15), Inches(3.8), Inches(0.38), 12.5, INK, True, align=PP_ALIGN.CENTER)
    slides_for_preview.append(("SOTA benchmark", ["Rank #12/35", "F1 0.949", "Mean SOTA 0.892"], [ROOT / "plots" / "plot11_paper_comparison.png"]))

    # 16. Conclusion
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, INK)
    add_text(s, "Final takeaways", Inches(0.75), Inches(0.62), Inches(6.3), Inches(0.5),
             34, WHITE, True, font=FONT_HEAD)
    takeaways = [
        ("Best result", "VLM-CLIP reaches macro F1=0.949."),
        ("Data claim", "Training uses field images, class-balanced augmentation, and generated symptom text; no Kaggle image corpus."),
        ("Privacy claim", "Leaf images, field logs, and RAG documents remain local; FedAvg exchanges weights."),
        ("Next work", "Add differential privacy, physically separated estates, more crops/diseases, and live IoT personalization."),
    ]
    for i, (head, body) in enumerate(takeaways):
        y = Inches(1.65 + i * 1.0)
        add_text(s, head, Inches(0.9), y, Inches(2.2), Inches(0.3), 16, GREEN2, True)
        add_text(s, body, Inches(3.0), y, Inches(8.2), Inches(0.45), 18, WHITE if i != 1 else RGBColor(255, 232, 205), True)
    add_band(s, Inches(0.75), Inches(6.7), Inches(11.6), Inches(0.04), GREEN2)
    add_text(s, "FarmFederate: field-realistic, multimodal, and privacy-preserving tea disease diagnosis.",
             Inches(0.8), Inches(6.9), Inches(10.6), Inches(0.3), 12.5, RGBColor(218, 232, 224))
    add_text(s, "16 / 16", Inches(12.05), Inches(7.12), Inches(0.65), Inches(0.18), 8, RGBColor(199, 213, 206), align=PP_ALIGN.RIGHT)
    slides_for_preview.append(("Final takeaways", [h + ": " + b for h, b in takeaways], []))

    prs.save(OUT)
    return slides_for_preview


def generate_previews(slide_specs):
    PREVIEW_DIR.mkdir(exist_ok=True)
    try:
        title_font = ImageFont.truetype("arialbd.ttf", 34)
        body_font = ImageFont.truetype("arial.ttf", 19)
        small_font = ImageFont.truetype("arial.ttf", 14)
    except OSError:
        title_font = body_font = small_font = ImageFont.load_default()

    thumbs = []
    for idx, (title, bullets, images) in enumerate(slide_specs, 1):
        img = Image.new("RGB", (960, 540), rgb_hex(PAPER))
        d = ImageDraw.Draw(img)
        d.rectangle([0, 0, 960, 18], fill=rgb_hex(GREEN2))
        d.text((44, 40), title, fill=rgb_hex(INK), font=title_font)
        y = 115
        for b in bullets[:5]:
            d.text((60, y), "- " + str(b)[:78], fill=rgb_hex(INK), font=body_font)
            y += 36
        if images:
            x0 = 540
            y0 = 130
            for p in images[:2]:
                p = Path(p)
                if p.exists():
                    with Image.open(p) as im:
                        im.thumbnail((350, 185))
                        img.paste(im.convert("RGB"), (x0, y0))
                    y0 += 205
        d.text((870, 505), f"{idx} / {len(slide_specs)}", fill=rgb_hex(GRAY), font=small_font)
        out = PREVIEW_DIR / f"slide_{idx:02d}.png"
        img.save(out)
        thumbs.append(out)

    montage = Image.new("RGB", (960, 1080), "white")
    for idx, p in enumerate(thumbs):
        with Image.open(p) as im:
            im.thumbnail((240, 135))
            x = (idx % 4) * 240
            y = (idx // 4) * 135
            montage.paste(im.convert("RGB"), (x, y))
    montage_path = PREVIEW_DIR / "montage.png"
    montage.save(montage_path)
    return montage_path


def inspect_pptx():
    prs = Presentation(OUT)
    slide_text_counts = []
    image_count = 0
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
            if shape.shape_type == 13:
                image_count += 1
        slide_text_counts.append(len(texts))
    return len(prs.slides), image_count, slide_text_counts


if __name__ == "__main__":
    specs = build_deck()
    montage = generate_previews(specs)
    slides, images, text_counts = inspect_pptx()
    print(f"updated={OUT}")
    print(f"backup={BACKUP}")
    print(f"slides={slides}")
    print(f"images={images}")
    print(f"text_counts={text_counts}")
    print(f"preview_montage={montage}")
