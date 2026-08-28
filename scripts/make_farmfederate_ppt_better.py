from __future__ import annotations

import datetime as dt
import shutil
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "FarmFederate_ppt_v2.pptx"
BACKUP = ROOT / f"FarmFederate_ppt_v2_before_better_update_{dt.datetime.now():%Y%m%d_%H%M%S}.pptx"
PREVIEW_DIR = ROOT / "ppt_previews_v2_better"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

INK = RGBColor(9, 14, 28)
INK2 = RGBColor(18, 29, 48)
PAPER = RGBColor(250, 249, 245)
PAPER2 = RGBColor(240, 248, 244)
GREEN = RGBColor(21, 111, 70)
GREEN2 = RGBColor(48, 148, 96)
TEAL = RGBColor(13, 116, 144)
BLUE = RGBColor(37, 99, 235)
RED = RGBColor(198, 40, 40)
AMBER = RGBColor(224, 138, 0)
PURPLE = RGBColor(109, 40, 217)
GRAY = RGBColor(98, 112, 126)
LIGHT = RGBColor(222, 232, 226)
WHITE = RGBColor(255, 255, 255)
MINT_TEXT = RGBColor(203, 225, 215)

FONT_HEAD = "Aptos Display"
FONT_BODY = "Aptos"


def set_fill(shape, color: RGBColor, transparency: int | None = None):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if transparency is not None:
        shape.fill.transparency = transparency


def no_line(shape):
    shape.line.fill.background()


def add_bg(slide, color=PAPER):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    set_fill(bg, color)
    no_line(bg)
    return bg


def add_text(slide, text, x, y, w, h, size=20, color=INK, bold=False,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, font=FONT_BODY,
             italic=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    return box


def add_bar(slide, x, y, w, h, color, transparency=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    set_fill(shp, color, transparency)
    no_line(shp)
    return shp


def add_round_rect(slide, x, y, w, h, color, radius=True, line=LIGHT):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h)
    set_fill(shp, color)
    if line:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    else:
        no_line(shp)
    return shp


def add_title(slide, idx, title, subtitle=None, section=None, dark=False):
    c_title = WHITE if dark else INK
    c_sub = MINT_TEXT if dark else GRAY
    c_sec = GREEN2 if dark else GREEN
    if section:
        add_text(slide, section.upper(), Inches(0.55), Inches(0.3), Inches(3.0), Inches(0.22),
                 8.5, c_sec, True)
    add_text(slide, title, Inches(0.55), Inches(0.52), Inches(8.8), Inches(0.58),
             28, c_title, True, font=FONT_HEAD)
    if subtitle:
        add_text(slide, subtitle, Inches(0.57), Inches(1.12), Inches(9.7), Inches(0.36),
                 12.5, c_sub)
    add_text(slide, f"{idx} / 18", Inches(12.05), Inches(7.12), Inches(0.7), Inches(0.18),
             8, c_sub, align=PP_ALIGN.RIGHT)
    add_text(slide, "FarmFederate", Inches(0.55), Inches(7.13), Inches(1.6), Inches(0.16),
             7.5, c_sub)


def add_metric(slide, value, label, x, y, w, color=GREEN, note=None, dark=False):
    add_text(slide, value, x, y, w, Inches(0.42), 24, color, True, align=PP_ALIGN.CENTER, font=FONT_HEAD)
    add_text(slide, label, x, y + Inches(0.48), w, Inches(0.3), 9.8,
             WHITE if dark else INK, True, align=PP_ALIGN.CENTER)
    if note:
        add_text(slide, note, x, y + Inches(0.78), w, Inches(0.25), 8.4,
                 MINT_TEXT if dark else GRAY, align=PP_ALIGN.CENTER)


def add_pill(slide, text, x, y, w, color=GREEN, text_color=WHITE):
    shp = add_round_rect(slide, x, y, w, Inches(0.3), color, line=None)
    add_text(slide, text, x + Inches(0.08), y + Inches(0.06), w - Inches(0.16), Inches(0.14),
             8.2, text_color, True, align=PP_ALIGN.CENTER)
    return shp


def add_image_contain(slide, img_path, x, y, w, h, border=True, bg=WHITE):
    img_path = Path(img_path)
    if border:
        add_round_rect(slide, x, y, w, h, bg, line=LIGHT)
    if not img_path.exists():
        add_text(slide, f"Missing: {img_path.name}", x + Inches(0.1), y + Inches(0.1),
                 w - Inches(0.2), Inches(0.4), 10, RED)
        return None
    with Image.open(img_path) as im:
        iw, ih = im.size
    scale = min(w / iw, h / ih)
    pw, ph = int(iw * scale), int(ih * scale)
    px = x + (w - pw) / 2
    py = y + (h - ph) / 2
    return slide.shapes.add_picture(str(img_path), px, py, width=pw, height=ph)


def add_table(slide, rows, x, y, w, h, widths=None, header_color=INK2, font_size=8.8):
    n_rows, n_cols = len(rows), len(rows[0])
    shape = slide.shapes.add_table(n_rows, n_cols, x, y, w, h)
    table = shape.table
    if widths:
        for i, frac in enumerate(widths):
            table.columns[i].width = int(w * frac)
    for r_i, row in enumerate(rows):
        for c_i, cell_val in enumerate(row):
            cell = table.cell(r_i, c_i)
            cell.text = str(cell_val)
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.02)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if c_i > 0 else PP_ALIGN.LEFT
            if not p.runs:
                p.add_run()
            run = p.runs[0]
            run.font.name = FONT_BODY
            run.font.size = Pt(font_size)
            run.font.bold = r_i == 0
            run.font.color.rgb = WHITE if r_i == 0 else INK
            cell.fill.solid()
            if r_i == 0:
                cell.fill.fore_color.rgb = header_color
            elif r_i % 2 == 0:
                cell.fill.fore_color.rgb = PAPER2
            else:
                cell.fill.fore_color.rgb = WHITE
    return shape


def add_callout(slide, heading, body, x, y, w, h, color=GREEN):
    add_bar(slide, x, y, Inches(0.07), h, color)
    add_text(slide, heading, x + Inches(0.2), y, w - Inches(0.2), Inches(0.25),
             12.5, color, True)
    add_text(slide, body, x + Inches(0.2), y + Inches(0.31), w - Inches(0.2),
             h - Inches(0.31), 10.5, INK)


def add_arrow(slide, x1, y1, x2, y2, color=GRAY):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(2)
    line.line.end_arrowhead = True


def add_flow(slide, text, x, y, w, h, color):
    shp = add_round_rect(slide, x, y, w, h, color, line=None)
    add_text(slide, text, x + Inches(0.1), y + Inches(0.11), w - Inches(0.2), h - Inches(0.16),
             10.5, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    return shp


def add_chart(slide, categories, series, x, y, w, h, chart_type=XL_CHART_TYPE.COLUMN_CLUSTERED,
              colors=None, legend=False):
    data = CategoryChartData()
    data.categories = categories
    for name, vals in series:
        data.add_series(name, vals)
    chart = slide.shapes.add_chart(chart_type, x, y, w, h, data).chart
    chart.has_legend = legend
    if legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.tick_labels.font.size = Pt(8)
    chart.category_axis.tick_labels.font.size = Pt(8)
    chart.plots[0].has_data_labels = True
    chart.plots[0].data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    chart.plots[0].data_labels.font.size = Pt(8)
    if colors:
        for si, ser in enumerate(chart.series):
            ser.format.fill.solid()
            ser.format.fill.fore_color.rgb = colors[si % len(colors)]
    return chart


def slide(prs, idx, title, subtitle=None, section=None, dark=False):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, INK if dark else PAPER)
    add_bar(s, 0, 0, SLIDE_W, Inches(0.15), GREEN2)
    add_title(s, idx, title, subtitle, section, dark)
    return s


def build_deck():
    if OUT.exists():
        shutil.copy2(OUT, BACKUP)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    preview_notes = []

    # 1 cover
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, INK)
    add_bar(s, 0, 0, SLIDE_W, Inches(0.16), GREEN2)
    add_bar(s, Inches(8.05), 0, Inches(0.18), SLIDE_H, GREEN2)
    add_image_contain(s, ROOT / "disease_annotated" / "_grid_overview.jpg",
                      Inches(8.55), Inches(0.55), Inches(4.15), Inches(5.95), border=False)
    add_text(s, "FarmFederate", Inches(0.75), Inches(0.85), Inches(6.5), Inches(0.65),
             42, WHITE, True, font=FONT_HEAD)
    add_text(s, "Multimodal Federated Learning for Tea Leaf Disease Detection",
             Inches(0.78), Inches(1.6), Inches(6.55), Inches(0.55), 17, MINT_TEXT)
    add_text(s, "Ayush Debnath", Inches(0.78), Inches(3.25), Inches(5.6), Inches(0.38),
             24, WHITE, True, font=FONT_HEAD)
    add_text(s, "22ME31010  ·  Mechanical Engineering",
             Inches(0.8), Inches(3.78), Inches(5.6), Inches(0.28), 15, MINT_TEXT, True)
    add_text(s, "Supervisor:  Prof. Sudip Misra",
             Inches(0.8), Inches(4.42), Inches(5.4), Inches(0.28), 14, WHITE, True)
    add_text(s, "Department of Computer Science & Engineering\nIndian Institute of Technology Kharagpur",
             Inches(0.8), Inches(4.86), Inches(6.2), Inches(0.55), 13, MINT_TEXT, True)
    for i, (v, lab, col) in enumerate([("5", "diseases", GREEN2), ("18", "models", TEAL),
                                       ("0.949", "best F1", AMBER), ("0", "Kaggle images", RED)]):
        add_metric(s, v, lab, Inches(0.75 + i * 1.55), Inches(6.0), Inches(1.25), col, dark=True)
    preview_notes.append(("Cover", ["Author/supervisor details", "5 diseases", "F1 0.949"], [ROOT / "disease_annotated" / "_grid_overview.jpg"]))

    # 2 agenda / contribution snapshot
    s = slide(prs, 2, "Paper-Aligned Defense Story", "The deck now carries the final paper structure, data claims, and result tables.", "Roadmap")
    add_table(s, [
        ["Block", "Slides", "What it proves"],
        ["Problem + data", "3-6", "Five tea diseases, real field source, no Kaggle image corpus"],
        ["Method", "7-9", "Local multimodal model, FedAvg, training safeguards"],
        ["Results", "10-15", "VLM top-line gain, federated retention, full ranking, SOTA context"],
        ["Deployment", "16-18", "RAG fallback, privacy story, limitations, future work"],
    ], Inches(0.75), Inches(1.65), Inches(7.0), Inches(2.3), widths=[0.22, 0.16, 0.62], font_size=9.5)
    add_callout(s, "Central claim",
                "FarmFederate gets the best performance from multimodal fusion while keeping raw leaf images, field logs, and advisory documents local to each estate.",
                Inches(8.1), Inches(1.75), Inches(4.3), Inches(1.25), GREEN)
    add_metric(s, "98.6%", "VLM federated retention", Inches(8.3), Inches(3.75), Inches(1.9), GREEN)
    add_metric(s, "0.038", "F1 over best ViT", Inches(10.3), Inches(3.75), Inches(1.75), TEAL)
    add_metric(s, "0.462", "F1 over best LLM", Inches(9.35), Inches(5.25), Inches(1.85), RED)
    preview_notes.append(("Roadmap", ["Problem/data", "Method", "Results", "Deployment"], []))

    # 3 five diseases
    s = slide(prs, 3, "Five Tea Leaf Diseases in Scope", "The paper now briefly describes what each class looks like in the field.", "Dataset")
    disease_rows = [
        ["Class", "Real images", "Field cue", "Type"],
        ["Leaf blight", "35", "Brown water-soaked necrosis", "Fungal/spotting"],
        ["Leaf hoppers", "8", "Puncture lesions and hopper damage", "Insect"],
        ["Leaf rust", "62", "Orange-yellow pustules", "Rust/fungal"],
        ["Looper caterpillars", "48", "Ragged holes, skeletonised patches", "Insect"],
        ["Mosquito bug", "47", "Corky wounds, shoot-tip dieback", "Insect"],
    ]
    add_table(s, disease_rows, Inches(0.65), Inches(1.55), Inches(7.4), Inches(3.1),
              widths=[0.26, 0.16, 0.42, 0.16], font_size=8.8)
    thumb_paths = [ROOT / "plots" / f for f in [
        "real_leaf_blight.jpg", "real_leaf_hoppers.jpg", "real_leaf_rust.jpg",
        "real_looper_caterpillars.jpg", "real_mosquito_bug.jpg"
    ]]
    for i, p in enumerate(thumb_paths):
        add_image_contain(s, p, Inches(8.45 + (i % 2) * 2.05), Inches(1.55 + (i // 2) * 1.53),
                          Inches(1.75), Inches(1.12), border=True)
    add_callout(s, "Why this matters",
                "The dataset mixes fungal spotting, rust pustules, and pest injury. The classifier must learn disease-region texture, not just broad leaf colour.",
                Inches(0.85), Inches(5.15), Inches(11.5), Inches(0.78), GREEN)
    preview_notes.append(("Five diseases", ["class counts", "field cues", "disease type"], thumb_paths[:2]))

    # 4 dataset overview with chart
    s = slide(prs, 4, "Dataset Numbers Used in the Paper", "These are the final values reflected in FarmFederate_Paper_v3.", "Dataset")
    add_table(s, [
        ["Data item", "Count", "Use in paper"],
        ["Field photos", "200", "Original localized image source"],
        ["YOLO-OBB boxes", "371", "Disease-region annotation"],
        ["Training-run image tensors", "792", "200 real + 592 minority synthetic fill"],
        ["Text descriptions", "3,000", "600 generated symptom records per class"],
        ["VLM train / val pairs", "633 / 79", "Paired text-image validation"],
    ], Inches(0.65), Inches(1.58), Inches(6.1), Inches(3.0), widths=[0.38, 0.18, 0.44], font_size=8.8)
    add_chart(s, ["Field\nphotos", "OBB\nboxes", "Train\nimages", "Text\nrows"],
              [("Count", [200, 371, 792, 3000])],
              Inches(7.25), Inches(1.65), Inches(5.2), Inches(3.1),
              colors=[GREEN])
    add_pill(s, "No external Kaggle image corpus in reported training", Inches(1.15), Inches(5.35), Inches(4.2), RED)
    add_text(s, "Colab may write Kaggle credentials, but the paper claim is based on field images + synthetic minority fill + generated symptom text.",
             Inches(5.6), Inches(5.32), Inches(6.3), Inches(0.45), 11, INK, True)
    preview_notes.append(("Dataset numbers", ["200 photos", "371 OBB", "792 train images", "3000 text"], []))

    # 5 class balance / training split
    s = slide(prs, 5, "Class Balance and Train/Validation Splits", "The deck separates the exported balanced folder from the actual reported training loader.", "Dataset")
    add_chart(s, ["Blight", "Hoppers", "Rust", "Looper", "Mosquito"],
              [("Real images", [35, 8, 62, 48, 47])],
              Inches(0.75), Inches(1.6), Inches(5.75), Inches(3.15), colors=[TEAL])
    add_table(s, [
        ["Loader", "Train", "Validation", "Note"],
        ["Text", "2,400", "300", "3,000 rows loaded"],
        ["Image", "633", "79", "792 tensors"],
        ["VLM pairs", "633", "79", "Label-matched pairs"],
    ], Inches(7.05), Inches(1.75), Inches(5.0), Inches(1.75), widths=[0.27, 0.18, 0.22, 0.33], font_size=9)
    add_table(s, [
        ["Text difficulty", "Share", "Purpose"],
        ["Clear templates", "8%", "Easy class-specific symptoms"],
        ["Slightly mixed", "22%", "Mild ambiguity"],
        ["Heavily mixed", "70%", "Forces image evidence in VLM"],
    ], Inches(7.05), Inches(4.15), Inches(5.0), Inches(1.75), widths=[0.36, 0.18, 0.46], font_size=9)
    preview_notes.append(("Class balance", ["real image class counts", "splits", "text difficulty"], []))

    # 6 pipeline
    s = slide(prs, 6, "Field-First Data Pipeline", "The pipeline in the deck now matches the Colab logs and paper wording.", "Pipeline")
    items = [
        ("Field collection\n200 images", GREEN),
        ("Sorted folders\n5 diseases", TEAL),
        ("YOLO-OBB labels\n371 boxes", BLUE),
        ("Minority fill\n+592 hoppers", AMBER),
        ("Generated text\n600/class", RED),
    ]
    x0 = 0.62
    for i, (txt, col) in enumerate(items):
        add_flow(s, txt, Inches(x0 + i * 2.45), Inches(2.0), Inches(1.9), Inches(0.85), col)
        if i < len(items) - 1:
            add_arrow(s, Inches(x0 + i * 2.45 + 1.9), Inches(2.43), Inches(x0 + (i + 1) * 2.45 - 0.08), Inches(2.43))
    add_image_contain(s, ROOT / "plots" / "plot21_dataset_comparison.png",
                      Inches(0.9), Inches(3.65), Inches(5.3), Inches(2.1), border=True)
    add_callout(s, "Paper-safe wording",
                "The balanced export folder is described as an inspection/download artifact. The reported training run uses the sorted real dataset and synthetic fill for the rare leaf-hoppers class.",
                Inches(6.75), Inches(3.7), Inches(5.2), Inches(1.5), GREEN)
    preview_notes.append(("Pipeline", ["field collection", "OBB labels", "minority fill", "generated text"], [ROOT / "plots" / "plot21_dataset_comparison.png"]))

    # 7 architecture
    s = slide(prs, 7, "System Architecture", "Local text, image, fusion, retrieval, and aggregation components are separated.", "Method")
    add_flow(s, "Symptom text\nLLM encoder", Inches(0.8), Inches(1.9), Inches(1.8), Inches(0.8), TEAL)
    add_flow(s, "Leaf image\nViT encoder", Inches(0.8), Inches(3.95), Inches(1.8), Inches(0.8), GREEN)
    add_flow(s, "8 VLM fusion\nstrategies", Inches(4.0), Inches(2.9), Inches(2.1), Inches(0.95), RED)
    add_flow(s, "Disease\nclassifier", Inches(7.15), Inches(2.9), Inches(1.8), Inches(0.95), BLUE)
    add_flow(s, "Local FAISS\nRAG advice", Inches(10.05), Inches(1.9), Inches(2.0), Inches(0.8), AMBER)
    add_flow(s, "FedAvg server\nweights only", Inches(10.05), Inches(3.95), Inches(2.0), Inches(0.8), INK2)
    for a in [(2.6, 2.3, 4.0, 3.15), (2.6, 4.35, 4.0, 3.6), (6.1, 3.38, 7.15, 3.38),
              (8.95, 3.15, 10.05, 2.3), (8.95, 3.6, 10.05, 4.35)]:
        add_arrow(s, *(Inches(v) for v in a))
    add_callout(s, "Privacy boundary",
                "Raw leaf images, symptom logs, and treatment documents remain on garden devices. The global server receives only model updates.",
                Inches(1.1), Inches(5.7), Inches(10.7), Inches(0.68), RED)
    preview_notes.append(("Architecture", ["LLM", "ViT", "8 VLM fusions", "FedAvg", "RAG"], []))

    # 8 training config
    s = slide(prs, 8, "Training Configuration and Anti-Collapse Controls", "The final paper reports the safeguards that keep all five classes active.", "Training")
    add_table(s, [
        ["Component", "Setting", "Why it matters"],
        ["Centralized epochs", "12", "Stable comparison across LLM/ViT/VLM"],
        ["Federated rounds", "8", "Non-IID privacy comparison"],
        ["Focal loss", "enabled", "Focuses learning on hard/rare classes"],
        ["Diversity loss", "weight=1.0", "Prevents one-class collapse"],
        ["Class weights", "sqrt damped, max 10x", "Avoids over-weight instability"],
        ["Hardware", "Tesla T4 GPU", "Colab-realistic training setup"],
    ], Inches(0.75), Inches(1.55), Inches(6.65), Inches(3.6), widths=[0.28, 0.24, 0.48], font_size=8.5)
    add_image_contain(s, ROOT / "plots" / "plot26_stress_distribution.png",
                      Inches(7.85), Inches(1.6), Inches(4.55), Inches(3.0), border=True)
    add_metric(s, "100%", "class diversity target", Inches(8.15), Inches(5.1), Inches(1.75), GREEN)
    add_metric(s, "AMP", "mixed precision", Inches(10.3), Inches(5.1), Inches(1.45), TEAL)
    preview_notes.append(("Training config", ["12 epochs", "8 rounds", "focal loss", "diversity loss"], [ROOT / "plots" / "plot26_stress_distribution.png"]))

    # 9 model families
    s = slide(prs, 9, "Model Families Compared Under One Protocol", "5 LLM + 5 ViT + 8 VLM fusion architectures, all reported in the paper.", "Models")
    add_table(s, [
        ["Family", "Count", "Best model", "Best F1", "Mean F1"],
        ["LLM", "5", "BERT-tiny", "0.487", "0.450"],
        ["ViT", "5", "EfficientNet", "0.911", "0.883"],
        ["VLM", "8", "CLIP fusion", "0.949", "0.915"],
    ], Inches(0.8), Inches(1.6), Inches(5.6), Inches(1.95), widths=[0.18, 0.14, 0.34, 0.17, 0.17], font_size=9.5)
    add_chart(s, ["LLM", "ViT", "VLM"], [("Best", [0.487, 0.911, 0.949]), ("Mean", [0.450, 0.883, 0.915])],
              Inches(6.95), Inches(1.55), Inches(5.3), Inches(3.05), colors=[GREEN, BLUE], legend=True)
    add_image_contain(s, ROOT / "plots" / "plot04_model_type_overview.png",
                      Inches(0.9), Inches(4.35), Inches(11.5), Inches(1.85), border=True)
    preview_notes.append(("Model families", ["LLM 5", "ViT 5", "VLM 8", "CLIP 0.949"], [ROOT / "plots" / "plot04_model_type_overview.png"]))

    # 10 full ranking
    s = slide(prs, 10, "Unified Ranking: All 18 Models", "A compact table now carries the actual ranking data, not only charts.", "Results")
    rank_rows = [
        ["Rank", "Model", "Type", "F1"],
        ["1", "VLM-CLIP", "VLM", "0.949"],
        ["2", "VLM-BLIP-2", "VLM", "0.937"],
        ["3", "VLM-attention", "VLM", "0.924"],
        ["4", "EfficientNet", "ViT", "0.911"],
        ["5", "VLM-flamingo", "VLM", "0.911"],
        ["6", "VLM-unified_io", "VLM", "0.911"],
        ["7", "ConvNeXT-tiny", "ViT", "0.899"],
        ["8", "VLM-gated", "VLM", "0.899"],
        ["9", "VLM-coca", "VLM", "0.899"],
        ["10", "VLM-concat", "VLM", "0.886"],
        ["11", "DeiT-tiny", "ViT", "0.873"],
        ["12", "Swin-tiny", "ViT", "0.873"],
        ["13", "ViT-Base", "ViT", "0.861"],
        ["14", "BERT-tiny", "LLM", "0.487"],
        ["18", "MobileBERT", "LLM", "0.417"],
    ]
    add_table(s, rank_rows[:9], Inches(0.65), Inches(1.45), Inches(5.9), Inches(4.8), widths=[0.13, 0.47, 0.20, 0.20], font_size=8.4)
    add_table(s, [rank_rows[0]] + rank_rows[9:], Inches(6.85), Inches(1.45), Inches(5.8), Inches(4.8), widths=[0.13, 0.47, 0.20, 0.20], font_size=8.4)
    add_text(s, "Takeaway: all VLM fusions beat all LLM-only models; top VLM also beats the best image-only model.",
             Inches(1.0), Inches(6.45), Inches(11.4), Inches(0.3), 12.5, INK, True, align=PP_ALIGN.CENTER)
    preview_notes.append(("Unified ranking", ["18 models", "VLM-CLIP #1", "MobileBERT #18"], []))

    # 11 VLM fusion
    s = slide(prs, 11, "VLM Fusion Details", "The model comparison now includes the eight multimodal fusion variants.", "Results")
    vlm_rows = [
        ["Fusion", "F1", "Params (M)", "Note"],
        ["CLIP", "0.949", "15.74", "Best overall"],
        ["BLIP-2", "0.937", "15.94", "Second"],
        ["Attention", "0.924", "15.87", "Strong fusion"],
        ["Flamingo", "0.911", "16.14", "Tied group"],
        ["Unified-IO", "0.911", "17.19", "Largest"],
        ["Gated", "0.899", "15.80", "Mid"],
        ["CoCa", "0.899", "16.20", "Mid"],
        ["Concat", "0.886", "15.61", "Baseline fusion"],
    ]
    add_table(s, vlm_rows, Inches(0.7), Inches(1.45), Inches(5.55), Inches(4.35), widths=[0.35, 0.17, 0.22, 0.26], font_size=8.2)
    add_image_contain(s, ROOT / "plots" / "plot03_vlm_fusion_comparison.png",
                      Inches(6.75), Inches(1.55), Inches(5.65), Inches(3.9), border=True)
    add_metric(s, "0.063", "VLM F1 range", Inches(7.05), Inches(5.75), Inches(1.7), TEAL)
    add_metric(s, "15.61-17.19M", "parameter span", Inches(9.4), Inches(5.75), Inches(2.35), GREEN)
    preview_notes.append(("VLM fusion", ["CLIP", "BLIP-2", "Attention", "Concat"], [ROOT / "plots" / "plot03_vlm_fusion_comparison.png"]))

    # 12 encoder baselines
    s = slide(prs, 12, "Encoder Baselines: Text vs Image vs Fusion", "Side-by-side plots and a data table make the modality gap clear.", "Results")
    add_image_contain(s, ROOT / "plots" / "plot01_llm_comparison.png", Inches(0.55), Inches(1.5), Inches(3.8), Inches(1.95), border=True)
    add_image_contain(s, ROOT / "plots" / "plot02_vit_comparison.png", Inches(4.75), Inches(1.5), Inches(3.8), Inches(1.95), border=True)
    add_image_contain(s, ROOT / "plots" / "plot03_vlm_fusion_comparison.png", Inches(8.95), Inches(1.5), Inches(3.8), Inches(1.95), border=True)
    add_table(s, [
        ["Best by modality", "Model", "F1", "Interpretation"],
        ["Text-only", "BERT-tiny", "0.487", "Symptom text is intentionally ambiguous"],
        ["Image-only", "EfficientNet", "0.911", "Visual disease cues dominate"],
        ["Multimodal", "VLM-CLIP", "0.949", "Text+image fusion adds top-line gain"],
    ], Inches(1.0), Inches(4.35), Inches(11.1), Inches(1.6), widths=[0.22, 0.24, 0.12, 0.42], font_size=9.2)
    preview_notes.append(("Baselines", ["BERT 0.487", "EfficientNet 0.911", "CLIP 0.949"], [ROOT / "plots" / "plot02_vit_comparison.png"]))

    # 13 federated
    s = slide(prs, 13, "Centralized vs Federated: Utility Retention", "The privacy-preserving setup keeps nearly all VLM performance.", "Federated")
    add_table(s, [
        ["Model", "Centralized F1", "Federated F1", "Retention", "Winner"],
        ["LLM", "0.427", "0.460", "107.7%", "Federated"],
        ["ViT", "0.886", "0.886", "100.0%", "Tie"],
        ["VLM", "0.873", "0.861", "98.6%", "Centralized"],
    ], Inches(0.7), Inches(1.48), Inches(5.8), Inches(1.75), widths=[0.18, 0.23, 0.23, 0.18, 0.18], font_size=8.9)
    add_image_contain(s, ROOT / "plots" / "plot05_centralized_vs_federated.png",
                      Inches(6.9), Inches(1.45), Inches(5.35), Inches(2.45), border=True)
    add_image_contain(s, ROOT / "plots" / "plot27_federated_convergence.png",
                      Inches(0.85), Inches(4.25), Inches(5.4), Inches(2.0), border=True)
    add_callout(s, "Interpretation",
                "The federated comparison measures whether the multimodal gain survives aggregation, not whether a centrally pooled public image dataset was added.",
                Inches(7.05), Inches(4.35), Inches(5.0), Inches(1.1), GREEN)
    preview_notes.append(("Federated", ["LLM 107.7%", "ViT 100%", "VLM 98.6%"], [ROOT / "plots" / "plot05_centralized_vs_federated.png"]))

    # 14 RAG and visual fallback
    s = slide(prs, 14, "RAG Advisory + Visual Diagnostic Fallback", "Classification is paired with local treatment retrieval and annotated visual references.", "RAG")
    add_image_contain(s, ROOT / "plots" / "rag_02_score_heatmap.png",
                      Inches(0.65), Inches(1.45), Inches(4.8), Inches(2.85), border=True)
    add_image_contain(s, ROOT / "disease_annotated" / "_grid_overview.jpg",
                      Inches(5.85), Inches(1.45), Inches(4.05), Inches(4.5), border=True)
    add_table(s, [
        ["RAG element", "Paper value"],
        ["Knowledge base", "5 disease treatment docs"],
        ["Retrieval", "Top-1 correct for all classes"],
        ["Fallback", "OBB annotated leaf reference"],
        ["Privacy", "Per-estate FAISS stays local"],
    ], Inches(10.25), Inches(1.65), Inches(2.55), Inches(2.25), widths=[0.45, 0.55], font_size=8.3)
    add_callout(s, "Operational value",
                "When symptom text is uncertain, FarmFederate gives a disease label, visual region evidence, and treatment advice without uploading raw records.",
                Inches(10.25), Inches(4.35), Inches(2.55), Inches(1.15), TEAL)
    preview_notes.append(("RAG fallback", ["Top-1 retrieval", "OBB grid", "Local FAISS"], [ROOT / "disease_annotated" / "_grid_overview.jpg"]))

    # 15 heatmap and ablation
    s = slide(prs, 15, "Ablation and Model Heatmap", "Results stay in the results section; conclusion now only concludes.", "Evidence")
    add_image_contain(s, ROOT / "plots" / "plot13_heatmap.png",
                      Inches(0.6), Inches(1.38), Inches(6.3), Inches(4.8), border=True)
    add_table(s, [
        ["Configuration", "Text", "Image", "F1"],
        ["Text-only BERT-tiny", "Yes", "No", "0.487"],
        ["Image-only EfficientNet", "No", "Yes", "0.911"],
        ["Multimodal VLM-CLIP", "Yes", "Yes", "0.949"],
    ], Inches(7.25), Inches(1.65), Inches(5.0), Inches(1.65), widths=[0.48, 0.14, 0.16, 0.22], font_size=8.8)
    add_image_contain(s, ROOT / "plots" / "plot10d_modality_contribution.png",
                      Inches(7.35), Inches(3.75), Inches(4.75), Inches(2.35), border=True)
    preview_notes.append(("Ablation", ["Text 0.487", "Image 0.911", "VLM 0.949"], [ROOT / "plots" / "plot13_heatmap.png"]))

    # 16 SOTA
    s = slide(prs, 16, "SOTA Context: Competitive Performance + Privacy", "FarmFederate is near the leading group while adding multimodal FL and local RAG.", "Benchmark")
    add_image_contain(s, ROOT / "plots" / "plot11_paper_comparison.png",
                      Inches(0.65), Inches(1.35), Inches(6.7), Inches(4.9), border=True)
    add_table(s, [
        ["Metric", "Value", "Meaning"],
        ["FarmFederate rank", "#12 / 35", "Top tier of compared tea/plant disease models"],
        ["FarmFederate F1", "0.949", "Best VLM-CLIP result"],
        ["Mean SOTA F1", "0.892", "FarmFederate exceeds benchmark mean"],
        ["Best literature F1", "0.995", "Highest image-centric tea CNN"],
    ], Inches(7.75), Inches(1.62), Inches(4.85), Inches(2.25), widths=[0.36, 0.22, 0.42], font_size=8.4)
    add_callout(s, "Fair positioning",
                "The paper presents FarmFederate as competitive, not universally best: its novelty is multimodal federated diagnosis with local advisory retrieval.",
                Inches(7.9), Inches(4.35), Inches(4.45), Inches(1.08), GREEN)
    preview_notes.append(("SOTA", ["rank #12/35", "F1 0.949", "mean SOTA 0.892"], [ROOT / "plots" / "plot11_paper_comparison.png"]))

    # 17 deployment, limitations, future
    s = slide(prs, 17, "Deployment Notes, Limitations, and Future Work", "The conclusion now combines limitations and future work compactly.", "Deployment")
    add_table(s, [
        ["Area", "Current paper", "Next step"],
        ["Data", "5 localized tea diseases", "Expand crop/disease coverage"],
        ["Privacy", "FedAvg, local records", "Formal differential privacy"],
        ["Clients", "Simulated non-IID clients", "Physically separated estates"],
        ["Models", "Colab/edge-sized encoders", "Compare larger foundation VLMs"],
        ["Context", "Weather/IoT adjustments", "Live sensor streams"],
    ], Inches(0.75), Inches(1.55), Inches(7.1), Inches(3.05), widths=[0.2, 0.38, 0.42], font_size=8.8)
    add_callout(s, "Deployment claim",
                "Images and symptom records stay local; ViT supports image-only diagnosis, text encoders support symptom logs, and CLIP-style fusion is used when both modalities are available.",
                Inches(8.25), Inches(1.68), Inches(4.1), Inches(1.55), GREEN)
    add_callout(s, "No-Kaggle guardrail",
                "The deck and paper state the same thing: no external Kaggle image corpus is used in the reported training pipeline.",
                Inches(8.25), Inches(4.05), Inches(4.1), Inches(1.0), RED)
    preview_notes.append(("Deployment", ["limitations", "future work", "no Kaggle"], []))

    # 18 final
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, INK)
    add_bar(s, 0, 0, SLIDE_W, Inches(0.16), GREEN2)
    add_text(s, "Final Takeaway", Inches(0.75), Inches(0.65), Inches(6.0), Inches(0.55),
             36, WHITE, True, font=FONT_HEAD)
    takeaways = [
        ("Performance", "VLM-CLIP is strongest: F1=0.949 across 18 models."),
        ("Data integrity", "Reported training uses field images, minority synthetic fill, and generated symptom text; no Kaggle image corpus."),
        ("Privacy", "FedAvg keeps raw images, logs, and advisory documents local."),
        ("Utility", "RAG + visual fallback turns classification into field-facing advice."),
    ]
    for i, (head, body) in enumerate(takeaways):
        y = Inches(1.65 + i * 1.0)
        add_text(s, head, Inches(0.95), y, Inches(2.0), Inches(0.27), 15, GREEN2, True)
        add_text(s, body, Inches(3.0), y, Inches(8.8), Inches(0.35), 18, WHITE if i != 1 else RGBColor(255, 232, 205), True)
    add_bar(s, Inches(0.75), Inches(6.58), Inches(11.8), Inches(0.04), GREEN2)
    add_text(s, "Ayush Debnath · 22ME31010 · Supervisor: Prof. Sudip Misra · IIT Kharagpur",
             Inches(0.82), Inches(6.85), Inches(9.8), Inches(0.25), 12, MINT_TEXT, True)
    add_text(s, "18 / 18", Inches(12.05), Inches(7.12), Inches(0.7), Inches(0.18), 8, MINT_TEXT, align=PP_ALIGN.RIGHT)
    preview_notes.append(("Final", [body for _, body in takeaways], []))

    prs.save(OUT)
    return preview_notes


def make_previews(notes):
    PREVIEW_DIR.mkdir(exist_ok=True)
    try:
        title_font = ImageFont.truetype("arialbd.ttf", 34)
        body_font = ImageFont.truetype("arial.ttf", 18)
        small_font = ImageFont.truetype("arial.ttf", 13)
    except OSError:
        title_font = body_font = small_font = ImageFont.load_default()
    thumbs = []
    for idx, (title, bullets, images) in enumerate(notes, 1):
        img = Image.new("RGB", (960, 540), "#faf9f5")
        d = ImageDraw.Draw(img)
        d.rectangle([0, 0, 960, 18], fill="#309460")
        d.text((42, 42), title, fill="#090e1c", font=title_font)
        y = 112
        for b in bullets[:5]:
            d.text((58, y), "- " + str(b)[:80], fill="#090e1c", font=body_font)
            y += 34
        x_img, y_img = 575, 115
        for p in images[:2]:
            p = Path(p)
            if p.exists():
                with Image.open(p) as im:
                    im.thumbnail((310, 170))
                    img.paste(im.convert("RGB"), (x_img, y_img))
                y_img += 188
        d.text((860, 506), f"{idx} / {len(notes)}", fill="#62707e", font=small_font)
        out = PREVIEW_DIR / f"slide_{idx:02d}.png"
        img.save(out)
        thumbs.append(out)

    montage = Image.new("RGB", (960, 675), "white")
    for idx, p in enumerate(thumbs):
        with Image.open(p) as im:
            im.thumbnail((192, 108))
            montage.paste(im.convert("RGB"), ((idx % 5) * 192, (idx // 5) * 112))
    montage_path = PREVIEW_DIR / "montage.png"
    montage.save(montage_path)
    return montage_path


def audit():
    prs = Presentation(OUT)
    text_blob = []
    image_count = 0
    table_count = 0
    chart_count = 0
    for s in prs.slides:
        for sh in s.shapes:
            if hasattr(sh, "text") and sh.text.strip():
                text_blob.append(sh.text)
            if sh.shape_type == 13:
                image_count += 1
            if sh.has_table:
                table_count += 1
            if sh.has_chart:
                chart_count += 1
    blob = "\n".join(text_blob).lower()
    checks = {
        "no kaggle": "no external kaggle image corpus" in blob or "no kaggle image corpus" in blob,
        "author": "ayush debnath" in blob and "22me31010" in blob,
        "supervisor": "prof. sudip misra" in blob,
        "best_f1": "0.949" in blob,
        "dataset_792": "792" in blob,
        "obb_371": "371" in blob,
        "rag": "rag" in blob,
        "fedavg": "fedavg" in blob,
    }
    return len(prs.slides), image_count, table_count, chart_count, checks


if __name__ == "__main__":
    notes = build_deck()
    montage = make_previews(notes)
    slides, images, tables, charts, checks = audit()
    print(f"updated={OUT}")
    print(f"backup={BACKUP}")
    print(f"slides={slides}")
    print(f"images={images}")
    print(f"tables={tables}")
    print(f"charts={charts}")
    print(f"checks={checks}")
    print(f"preview_montage={montage}")
